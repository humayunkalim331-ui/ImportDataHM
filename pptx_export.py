"""
Director-review PPTX export — built with python-pptx (the original tool hand-built raw
OOXML via JSZip because no pptx library was available in the browser; here we have one).
Same slide set and narration logic, generated a much simpler way.
"""
import io
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from aggregation import (
    num, aggregate_by, dominant_currency_prices, weighted_average, monthly_pivot,
    full_company_month_matrix, landed_cost_ranked, compute_rank, month_label,
    fmt_qty_mt, fmt_price_per_mt, detect_landed_cost_columns, landed_cost_formula_text,
)
from matching import company_key
from insights import build_series, build_insights_bullets
from charts import trend_png_bytes

NAVY = RGBColor(0x1F, 0x38, 0x64)
GOLD = RGBColor(0xC9, 0x9A, 0x4A)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xF5, 0xF7, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _text(slide, x, y, w, h, text, size=13, bold=False, color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def _header(slide, title, subtitle=None):
    _rect(slide, 0, 0, SLIDE_W, Inches(0.7), NAVY)
    _rect(slide, 0, Inches(0.7), SLIDE_W, Inches(0.035), GOLD)
    _text(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.5), title, size=22, bold=True, color=WHITE)
    if subtitle:
        _text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4), subtitle, size=12, color=GREY)


def _table(slide, x, y, w, headers, rows, row_h=Inches(0.42), font_size=11):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_h = row_h * n_rows
    gframe = slide.shapes.add_table(n_rows, n_cols, x, y, w, total_h)
    table = gframe.table
    col_w = int(w / n_cols)
    for c in range(n_cols):
        table.columns[c].width = Emu(col_w)
    for r in range(n_rows):
        table.rows[r].height = row_h

    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(htext)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = WHITE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_LIGHT if r % 2 == 1 else ROW_ALT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            p.font.size = Pt(font_size - 1)
            p.font.color.rgb = DARK
    return total_h


def _pptx_top_rows(rows, mapping, name_role, country_role):
    if not mapping or not mapping.get(name_role) or not mapping.get("qty"):
        return []
    agg = aggregate_by(rows, mapping[name_role], mapping, company_key)[:5]
    out = []
    for name, v in agg:
        country_str = ""
        if country_role and mapping.get(country_role):
            vals = v["values"].get(country_role, [])
            country_str = ", ".join(str(x) for x in vals[:3]) + (f" +{len(vals)-3} more" if len(vals) > 3 else "")
        out.append([
            name, country_str,
            fmt_qty_mt(v['qty']) if v["qty"] else "—",
            fmt_price_per_mt(v["priceAvg"], v["currency"]) if v["priceAvg"] is not None else "—",
        ])
    return out


def _pptx_region_rows(rows, mapping, region_role):
    if not mapping or not mapping.get(region_role) or not mapping.get("qty"):
        return []
    agg = aggregate_by(rows, mapping[region_role], mapping)[:6]
    return [[region, fmt_qty_mt(v['qty']) if v["qty"] else "—",
              fmt_price_per_mt(v["priceAvg"], v["currency"]) if v["priceAvg"] is not None else "—"]
             for region, v in agg]


def _pptx_full_matrix_rows(rows, mapping, company_role, max_companies, label):
    matrix = full_company_month_matrix(rows, mapping, company_role, max_companies)
    if not matrix or not matrix["companies"]:
        return None
    months = matrix["months"][-6:]
    headers = [label] + [month_label(mk) for mk in months]
    body_rows = []
    for name in matrix["companies"]:
        row_out = [name]
        for mk in months:
            cell = matrix["cell"].get(name, {}).get(mk)
            if not cell:
                row_out.append("—")
            else:
                avg = cell["wsum"] / cell["wqty"] if cell["wqty"] > 0 else (cell["sum"] / cell["n"] if cell["n"] else 0)
                row_out.append(f"{cell['qty']:,.0f} @ {avg:.2f}")
        body_rows.append(row_out)
    return {"headers": headers, "rows": body_rows, "rest": matrix["rest"], "total_companies": matrix["total_companies"]}


def build_pptx(item, manual, imp_rows, imp_mapping, has_import_file, imp_headers,
                exp_rows, exp_mapping, has_export_file,
                wits_rows, wits_mapping, data_view, period_label):
    """Returns BytesIO containing the .pptx file."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    show_import = data_view.get("import", True) and bool(imp_rows) and imp_mapping and imp_mapping.get("price")
    show_export = data_view.get("export", True) and bool(exp_rows) and exp_mapping and exp_mapping.get("price")

    # Slide 1 — Title
    s = _blank_slide(prs)
    _rect(s, 0, 0, SLIDE_W, Inches(0.18), NAVY)
    _rect(s, 0, Inches(0.18), SLIDE_W, Inches(0.04), GOLD)
    _text(s, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.1), item["name"], size=36, bold=True, color=DARK)
    _text(s, Inches(0.8), Inches(3.55), Inches(11.5), Inches(0.5), f"HS {item['hs']}  ·  {item['plant']}  ·  Origin {item.get('origin','—')}", size=15, color=GREY)
    _text(s, Inches(0.8), Inches(4.05), Inches(11.5), Inches(0.5), "Procurement Price Analysis — Director Review", size=15, bold=True, color=NAVY)
    subtitle = f"This report benchmarks our current price and supplier position for {item['name']} against matched import{' and export' if show_export else ''} customs data, and closes with a recommended next step."
    _text(s, Inches(0.8), Inches(4.6), Inches(10.8), Inches(1.3), subtitle, size=12, color=GREY)
    _text(s, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.4),
          f"Gatronova — RM Procurement Analyzer  ·  Generated {datetime.now().strftime('%m/%d/%Y')}  ·  Period: {period_label}",
          size=10, color=GREY)

    # Slide 2 — Our current position
    our_price = num(manual.get("price"))
    rank_line = "n/a — enter unit price and map the import price column"
    rank = compute_rank(imp_rows, imp_mapping, our_price) if imp_mapping else None
    if rank:
        rank_line = f"#{rank['rank']} of {rank['total']} matched importers{' in ' + rank['currency'] if rank['currency'] else ''} (lowest price = #1)"
    trend_imp = build_series(imp_rows, imp_mapping)
    if len(trend_imp) >= 2:
        trend_dir = "Rising" if trend_imp[-1]["price"] > trend_imp[0]["price"] else ("Falling" if trend_imp[-1]["price"] < trend_imp[0]["price"] else "Flat")
    else:
        trend_dir = "Not enough data points"
    vs_market_line = "n/a"
    imp_avg = None
    if imp_mapping and imp_mapping.get("price") and imp_rows and our_price is not None:
        prices, pairs, vs_cur, _ = dominant_currency_prices(imp_rows, imp_mapping)
        if prices:
            imp_avg = weighted_average(pairs)
            diff_pct = (our_price - imp_avg) / imp_avg * 100
            vs_market_line = f"{abs(diff_pct):.1f}% {'above' if diff_pct >= 0 else 'below'} quantity-weighted market average ({imp_avg:.2f}{' ' + vs_cur if vs_cur else ''})"

    scorecard_rows = [
        ["Quantity", str(manual.get("qty") or "n/a")],
        ["Unit Price", f"{manual.get('price') or 'n/a'} {manual.get('currency') or ''}".strip()],
        ["Supplier", str(manual.get("supplier") or "n/a")],
        ["Rank vs. matched importers", rank_line],
        ["Vs. market average", vs_market_line],
        ["Import price trend (selected period)", trend_dir],
    ]
    position_narration = (
        f"We are currently {vs_market_line}, ranked {rank_line}. The import price trend for the selected period is "
        f"{trend_dir.lower()}, which should factor into the timing of the next purchase."
        if our_price is not None and vs_market_line != "n/a" else
        "Enter our current quantity, unit price, and supplier above to populate the market ranking and comparison for this item."
    )
    s = _blank_slide(prs)
    _header(s, "Our Current Position")
    th = _table(s, Inches(0.8), Inches(1.2), Inches(9.0), ["Metric", "Value"], scorecard_rows)
    _text(s, Inches(0.8), Inches(1.2) + th + Inches(0.2), Inches(10.9), Inches(0.9), position_narration, size=12, color=GREY)

    # Slide 3 — Last importer data: top 5 competitors (month pivot)
    if show_import:
        pivot = monthly_pivot(imp_rows, imp_mapping, "importer", 5)
        if pivot:
            s = _blank_slide(prs)
            _header(s, "Last Importer Data — Top 5 Competitors (Import Customs Data)")
            tw = min(Inches(12.1), len(pivot["headers"]) * Inches(2.0))
            th = _table(s, Inches(0.6), Inches(1.1), tw, pivot["headers"], pivot["rows"])
            narration = (
                f"{pivot['cheapest']['name']} has offered the most competitive matched price in this window, at "
                f"{pivot['cheapest']['price']:.2f}{' ' + pivot['cheapest']['currency'] if pivot['cheapest']['currency'] else ''}. "
                "Use this table to benchmark upcoming negotiations against actual recent shipment prices rather than a blended average."
                if pivot.get("cheapest") else
                "This table lists actual shipment prices by month for our top 5 matched competitors by volume."
            )
            _text(s, Inches(0.6), Inches(1.1) + th + Inches(0.15), Inches(11.7), Inches(1.0), narration, size=11, color=GREY)

    # Slide 4 — Top importers
    if show_import:
        rows = _pptx_top_rows(imp_rows, imp_mapping, "importer", "origin")
        if rows:
            imp_agg = aggregate_by(imp_rows, imp_mapping["importer"], imp_mapping, company_key)
            leader = imp_agg[0]
            priced = [x for x in imp_agg if x[1]["priceAvg"] is not None]
            cheapest = min(priced, key=lambda x: x[1]["priceAvg"]) if priced else None
            narration = f"{leader[0]} is the largest matched importer by volume ({fmt_qty_mt(leader[1]['qty'])} MT)"
            if cheapest and cheapest[0] != leader[0]:
                narration += f", while {cheapest[0]} offers the most competitive matched price at {fmt_price_per_mt(cheapest[1]['priceAvg'], cheapest[1]['currency'])}"
            narration += ". Compare these terms against our current supplier before the next order."
            s = _blank_slide(prs)
            _header(s, f"Top Importers — HS {item['hs']}")
            th = _table(s, Inches(0.6), Inches(1.2), Inches(11.5), ["Importer", "Country", "Qty (MT)", "Unit Price (per MT)"], rows)
            _text(s, Inches(0.6), Inches(1.2) + th + Inches(0.2), Inches(11.7), Inches(0.9), narration, size=12, color=GREY)

    # Slide 5 — Region-wise: imports
    if show_import:
        rows = _pptx_region_rows(imp_rows, imp_mapping, "origin")
        if rows:
            region_agg = aggregate_by(imp_rows, imp_mapping["origin"], imp_mapping)
            top_region = region_agg[0]
            avg_str = fmt_price_per_mt(top_region[1]['priceAvg'], top_region[1]['currency']) if top_region[1]["priceAvg"] is not None else "n/a"
            narration = (f"{top_region[0]} is the dominant sourcing origin by volume for matched imports "
                         f"({fmt_qty_mt(top_region[1]['qty'])} MT), averaging {avg_str}.")
            narration += " Consider whether diversifying sourcing origins could reduce single-region dependency risk."
            s = _blank_slide(prs)
            _header(s, "Region-wise Comparison — Imports (by Country of Origin)")
            th = _table(s, Inches(0.6), Inches(1.2), Inches(10.5), ["Country of Origin", "Qty (MT)", "Avg Unit Price (per MT)"], rows)
            _text(s, Inches(0.6), Inches(1.2) + th + Inches(0.2), Inches(11.7), Inches(0.9), narration, size=12, color=GREY)

    # Slide 6 — Last exporter data: top 5 competitors
    if show_export:
        pivot = monthly_pivot(exp_rows, exp_mapping, "exporter", 5)
        if pivot:
            s = _blank_slide(prs)
            _header(s, "Last Exporter Data — Top 5 Competitors (Exporter Customs Data)")
            tw = min(Inches(12.1), len(pivot["headers"]) * Inches(2.0))
            th = _table(s, Inches(0.6), Inches(1.1), tw, pivot["headers"], pivot["rows"])
            narration = (
                f"{pivot['cheapest']['name']} has offered the most competitive matched export price in this window, "
                f"at {pivot['cheapest']['price']:.2f}{' ' + pivot['cheapest']['currency'] if pivot['cheapest']['currency'] else ''}."
                if pivot.get("cheapest") else
                "This table lists actual shipment prices by month for our top 5 matched exporters by volume."
            )
            _text(s, Inches(0.6), Inches(1.1) + th + Inches(0.15), Inches(11.7), Inches(0.9), narration, size=11, color=GREY)

    # Slide 7 — Top exporters
    if show_export:
        rows = _pptx_top_rows(exp_rows, exp_mapping, "exporter", "expcountry")
        if rows:
            exp_agg = aggregate_by(exp_rows, exp_mapping["exporter"], exp_mapping, company_key)
            leader = exp_agg[0]
            priced = [x for x in exp_agg if x[1]["priceAvg"] is not None]
            cheapest = min(priced, key=lambda x: x[1]["priceAvg"]) if priced else None
            narration = f"{leader[0]} is the largest matched exporter by volume ({fmt_qty_mt(leader[1]['qty'])} MT)"
            if cheapest and cheapest[0] != leader[0]:
                narration += f", while {cheapest[0]} offers the most competitive matched price at {fmt_price_per_mt(cheapest[1]['priceAvg'], cheapest[1]['currency'])}"
            narration += "."
            s = _blank_slide(prs)
            _header(s, f"Top Exporters — HS {item['hs']}")
            th = _table(s, Inches(0.6), Inches(1.2), Inches(11.5), ["Exporter", "Country", "Qty (MT)", "Unit Price (per MT)"], rows)
            _text(s, Inches(0.6), Inches(1.2) + th + Inches(0.2), Inches(11.7), Inches(0.9), narration, size=12, color=GREY)

    # Slide 8 — Region-wise: exports
    if show_export:
        rows = _pptx_region_rows(exp_rows, exp_mapping, "expcountry")
        if rows:
            region_agg = aggregate_by(exp_rows, exp_mapping["expcountry"], exp_mapping)
            top_region = region_agg[0]
            avg_str = fmt_price_per_mt(top_region[1]['priceAvg']) if top_region[1]["priceAvg"] is not None else "n/a"
            narration = (f"{top_region[0]} is the dominant exporting country by volume in matched records "
                         f"({fmt_qty_mt(top_region[1]['qty'])} MT), averaging {avg_str}"
                         f"{' ' + top_region[1]['currency'] if top_region[1]['currency'] else ''}.")
            s = _blank_slide(prs)
            _header(s, "Region-wise Comparison — Exports (by Exporting Country)")
            th = _table(s, Inches(0.6), Inches(1.2), Inches(10.5), ["Exporting Country", "Qty (MT)", "Avg Unit Price (per MT)"], rows)
            _text(s, Inches(0.6), Inches(1.2) + th + Inches(0.2), Inches(11.7), Inches(0.9), narration, size=12, color=GREY)

    # Slide 9 — Exporter data based on selected item (full month x supplier matrix)
    if show_import:
        matrix = _pptx_full_matrix_rows(imp_rows, imp_mapping, "supplier", 8, "Supplier (Exporter)")
        if matrix:
            rest_note = f" Showing top {len(matrix['headers'])-1} of {matrix['total_companies']} suppliers by total quantity." if matrix["rest"] > 0 else ""
            s = _blank_slide(prs)
            _header(s, "Exporter Data Based on Selected Item")
            tw = min(Inches(12.1), len(matrix["headers"]) * Inches(2.1))
            th = _table(s, Inches(0.6), Inches(1.1), tw, matrix["headers"], matrix["rows"], font_size=10)
            _text(s, Inches(0.6), Inches(1.1) + th + Inches(0.15), Inches(11.7), Inches(0.9),
                  f'Each cell: total quantity @ quantity-weighted average unit price, by month. From the "Supplier Name" '
                  f'column — the foreign seller is the exporter from our side as the buyer/importer.{rest_note}',
                  size=10, color=GREY)

    # Slide 10 — Importer data based on selected item (full month x importer matrix)
    if show_import:
        matrix = _pptx_full_matrix_rows(imp_rows, imp_mapping, "importer", 8, "Importer")
        if matrix:
            rest_note = f" Showing top {len(matrix['headers'])-1} of {matrix['total_companies']} importers by total quantity." if matrix["rest"] > 0 else ""
            s = _blank_slide(prs)
            _header(s, "Importer Data Based on Selected Item")
            tw = min(Inches(12.1), len(matrix["headers"]) * Inches(2.1))
            th = _table(s, Inches(0.6), Inches(1.1), tw, matrix["headers"], matrix["rows"], font_size=10)
            _text(s, Inches(0.6), Inches(1.1) + th + Inches(0.15), Inches(11.7), Inches(0.9),
                  f'Each cell: total quantity @ quantity-weighted average unit price, by month. From the '
                  f'"Importer Name on GD" column.{rest_note}', size=10, color=GREY)

    # Slide 11 — Landed cost comparison
    if show_import:
        ranked = landed_cost_ranked(imp_rows, imp_mapping, imp_headers, "supplier", limit=8)
        if ranked:
            rows = [[str(i + 1), name, fmt_qty_mt(qty), fmt_price_per_mt(lc)] for i, (name, lc, qty) in enumerate(ranked)]
            cheapest_name, cheapest_lc, _ = ranked[0]
            priciest_name, priciest_lc, _ = ranked[-1]
            gap_pct = ((priciest_lc - cheapest_lc) / cheapest_lc * 100) if cheapest_lc > 0 else None
            narration = f"{cheapest_name} is the most cost-effective supplier once real paid duties are included, at {fmt_price_per_mt(cheapest_lc)} PKR/MT."
            if gap_pct is not None and priciest_name != cheapest_name:
                narration += f" That's {gap_pct:.1f}% cheaper landed than {priciest_name}, the most expensive supplier shown here — a gap that unit price alone would not reveal."
            lc_cols = detect_landed_cost_columns(imp_headers)
            formula, note = landed_cost_formula_text(lc_cols)
            narration += f" Formula: {formula}.{note}"
            s = _blank_slide(prs)
            _header(s, "Landed Cost Comparison — Real Cost After Duties & Taxes")
            th = _table(s, Inches(0.6), Inches(1.2), Inches(10.5), ["#", "Supplier (Exporter)", "Qty (MT)", "Landed Cost / MT (PKR)"], rows)
            _text(s, Inches(0.6), Inches(1.2) + th + Inches(0.2), Inches(11.7), Inches(1.3), narration, size=10.5, color=GREY)

    # Slide 12 — Price trend
    trend_png = trend_png_bytes(imp_rows, imp_mapping, exp_rows, exp_mapping, wits_rows, wits_mapping, manual)
    if trend_png:
        trend_verb = "risen" if trend_dir == "Rising" else "fallen" if trend_dir == "Falling" else "stayed flat"
        if len(trend_imp) >= 2:
            trend_narration = (f"Import price has {trend_verb} over the selected period, moving from "
                               f"{trend_imp[0]['price']:.2f} to {trend_imp[-1]['price']:.2f}."
                               + (" The dashed line marks our current price for direct comparison against market movement." if our_price is not None else ""))
        else:
            trend_narration = "Not enough dated price points in the selected period to establish a trend — this chart will fill in as more dated shipments are matched."
        s = _blank_slide(prs)
        _header(s, "Price Trend")
        s.shapes.add_picture(io.BytesIO(trend_png), Inches(0.6), Inches(1.1), width=Inches(12.1), height=Inches(4.4))
        _text(s, Inches(0.6), Inches(5.65), Inches(11.7), Inches(0.8), trend_narration, size=12, color=GREY)

    # Slide 13 — Insights & recommendations
    bullets = build_insights_bullets(item, manual, imp_rows, imp_mapping, has_import_file,
                                      exp_rows, exp_mapping, has_export_file)
    s = _blank_slide(prs)
    _header(s, "Insights & Recommendations")
    _text(s, Inches(0.6), Inches(1.0), Inches(11.7), Inches(0.5),
          "The following points summarize our procurement position and recommended next steps, based on the matched market data above.",
          size=13, color=GREY)
    _text(s, Inches(0.6), Inches(1.6), Inches(11.7), Inches(4.9),
          "\n\n".join(f"•  {b}" for b in bullets), size=16, color=DARK)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
